from pptx import Presentation
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()

# Slide 1: Title Slide
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide_1.shapes.title
title.text = "The Art of Simplifying Life with Al"
subtitle = slide_1.placeholders[1]
subtitle.text = "Your Name"

# Slide 2: Introduction
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Introduction"
content = slide_2.shapes.add_textbox(left=100, top=100, width=600, height=300)
text_frame = content.text_frame
text_frame.text = "Welcome to the world of simplification with Al!"

# Slide 3: Benefits of Simplification
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_3.shapes.title
title.text = "Benefits of Simplification"
content = slide_3.shapes.add_textbox(left=100, top=100, width=600, height=300)
text_frame = content.text_frame
text_frame.text = "Why Simplify?"
p = text_frame.add_paragraph()
p.text = "Simplifying life can reduce stress, improve focus, and enhance productivity."

# Slide 4: Define Simplification
slide_4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_4.shapes.title
title.text = "Define Simplification"
content = slide_4.shapes.add_textbox(left=100, top=100, width=600, height=300)
text_frame = content.text_frame
text_frame.text = "What is Simplification?"
p = text_frame.add_paragraph()
p.text = "Simplification is the process of removing unnecessary complexity from tasks and routines."

# Add more content slides as needed

# Slide 9: Simplification Tips
slide_9 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_9.shapes.title
title.text = "Simplification Tips"
content = slide_9.shapes.add_textbox(left=100, top=100, width=600, height=300)
text_frame = content.text_frame
text_frame.text = "Practical Tips"
p = text_frame.add_paragraph()
p.text = "Here are some practical tips for simplifying your life with Al:"

# Slide 10: Conclusion
slide_10 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_10.shapes.title
title.text = "Conclusion"
content = slide_10.shapes.add_textbox(left=100, top=100, width=600, height=300)
text_frame = content.text_frame
text_frame.text = "Simplify and Thrive"
p = text_frame.add_paragraph()
p.text = "Embrace the art of simplification with Al and unlock a more peaceful and productive life!"

# Save the presentation
presentation.save("Simplifying_Life_With_Al.pptx")
